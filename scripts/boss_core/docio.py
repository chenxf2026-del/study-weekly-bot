"""boss_core.docio — 文档读取 / 截断 / 分节 (M0.1b, 从 run_pipeline_local 纯搬移)。

无状态工具, 只依赖 stdlib + lazy 文档库 (pypdf/python-docx/python-pptx) + boss_core.errors。
run_pipeline_local 顶部 re-export 这三个名字, 生产 import 方 (feishu_events / review_batch /
reviewed_docs_filter / meeting_summary_pipeline) 与测试零改动。
"""

from __future__ import annotations

import re
from pathlib import Path
from typing import Optional

from boss_core.errors import PipelineError


def _read_doc_text(path: Path) -> str:
    """读 review doc 文本。.md/.txt 直读; .pdf/.docx 提取文字 (飞书机器人收的多是这两种)。
    .doc (老 Word 二进制) 不支持, 明确引导转 .docx/.pdf。提取依赖缺失时报清晰错误。"""
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        try:
            from pypdf import PdfReader
        except ImportError as e:
            raise PipelineError("PDF 解析需 pypdf (pip install pypdf)") from e
        reader = PdfReader(str(path))
        return "\n\n".join((pg.extract_text() or "") for pg in reader.pages)
    if suffix == ".docx":
        try:
            import docx  # python-docx
        except ImportError as e:
            raise PipelineError("Word 解析需 python-docx (pip install python-docx)") from e
        d = docx.Document(str(path))
        parts = [p.text for p in d.paragraphs]
        for tbl in d.tables:                       # 表格文字也提取 (纪要/方案常含表)
            for row in tbl.rows:
                parts.append("\t".join(c.text for c in row.cells))
        return "\n".join(parts)
    if suffix == ".pptx":
        try:
            from pptx import Presentation  # python-pptx
        except ImportError as e:
            raise PipelineError("PPT 解析需 python-pptx (pip install python-pptx)") from e
        prs = Presentation(str(path))
        parts: list[str] = []
        for i, slide in enumerate(prs.slides, 1):
            parts.append(f"# 第 {i} 页")
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    parts.append(shape.text_frame.text)
                if shape.has_table:                       # 表格文字 (方案页常含表)
                    for row in shape.table.rows:
                        parts.append("\t".join(c.text for c in row.cells))
            notes = slide.notes_slide if slide.has_notes_slide else None
            if notes and notes.notes_text_frame and notes.notes_text_frame.text.strip():
                parts.append(f"[演讲备注] {notes.notes_text_frame.text}")
        return "\n".join(parts)
    if suffix in (".doc", ".ppt"):
        kind = "Word" if suffix == ".doc" else "PowerPoint"
        raise PipelineError(f"{suffix} (老 {kind} 二进制) 暂不支持, 请另存为 .docx/.pptx 或导出 PDF 再发")
    # .md / .txt / 其他: 当 UTF-8 文本读 (errors=replace 防个别坏字节炸掉)
    return path.read_text(encoding="utf-8", errors="replace")


def _cap_doc_text(content: str, cap: int) -> tuple[str, bool]:
    """文档正文超上限时智能截断: 保留前 70% + 后 30% (汇报稿结论常在末尾), 中间插标记。

    返回 (可能截断后的文本, 是否截断)。cap<=0 或未超限则原样返回。
    截断只影响塞进 Phase 0 parse prompt 的量, 不改磁盘上的原文档。
    """
    if cap <= 0 or len(content) <= cap:
        return content, False
    head_n = int(cap * 0.7)
    tail_n = cap - head_n
    marker = (
        f"\n\n[⚠ 文档过长 (原 {len(content)} 字符), 已截断为前 {head_n} + 后 {tail_n} 字符供解析; "
        f"中间省略约 {len(content) - cap} 字符。如需全文评审请拆分文档或调高 REVIEW_DOC_MAX_CHARS。]\n\n"
    )
    return content[:head_n] + marker + content[-tail_n:], True


def _extract_section_body(report_path: Path, heading_candidates: list[str]) -> Optional[str]:
    """读 report.md, 找第一个命中 heading_candidates 的 ### heading,
    返回该 heading 之后到下一个 ## / ### 之前的内容 (str). 找不到返回 None.
    Heading 用 prefix match 兼容 `### 共识 (5 评委都点头)` 等带后缀形式.
    """
    if not report_path.exists():
        return None
    try:
        text = report_path.read_text(encoding="utf-8")
    except (OSError, UnicodeDecodeError):
        return None
    for heading in heading_candidates:
        pattern = re.compile(rf"^{re.escape(heading)}[ \t]*[^\n]*$", re.MULTILINE)
        m = pattern.search(text)
        if not m:
            continue
        start = m.end()
        rest = text[start:]
        next_m = re.search(r"^#{2,3}[ \t]+", rest, re.MULTILINE)
        body = rest[:next_m.start()] if next_m else rest
        return body.strip()
    return None
